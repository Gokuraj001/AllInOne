from flask import Flask, render_template, request, send_file, current_app, redirect, url_for
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
import img2pdf
import pdfplumber
import os
import zipfile
import pymupdf as fitz  # PyMuPDF
import uuid
import pandas as pd
import tabula
import comtypes.client
import tempfile
import platform
import pythoncom
import win32com.client
import qrcode
from flask import jsonify
from pdf2docx import Converter
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from PIL import Image, ImageOps
from werkzeug.utils import secure_filename, send_from_directory
from docx import Document
from openpyxl import load_workbook


# ===================== APP SETUP =====================
def create_app():
    app = Flask(__name__)

    app.config.from_mapping(
        UPLOAD_FOLDER=os.path.join(app.root_path, 'uploads'),
        OUTPUT_FOLDER=os.path.join(app.root_path, 'outputs')
    )

    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
    os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

    

    # ===================== HELPER FUNCTION =====================
    def save_uploaded_file(file, file_type='pdf'):
        if not file or file.filename == '':
            return None, None

        filename = secure_filename(file.filename)

        allowed = {
            'pdf': ('.pdf',),
            'image': ('.jpg', '.jpeg', '.png')
        }

        if file_type in allowed:
            if not filename.lower().endswith(allowed[file_type]):
                return None, None

        filepath = os.path.join(current_app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        return filename, filepath
    
    ALLOWED_OFFICE_EXTENSIONS = {'doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx'}

    def allowed_office_file(filename):
        return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_OFFICE_EXTENSIONS
    
        # ===================== PDF CONVERTER HELPERS =====================
    def convert_pdf_to_word(pdf_path):
        output_filename = f"converted_{uuid.uuid4().hex}.docx"
        output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], output_filename)

        cv = Converter(pdf_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()

        return output_filename, output_path


    def convert_pdf_to_excel(pdf_path):
        output_filename = f"converted_{uuid.uuid4().hex}.xlsx"
        output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], output_filename)

        try:
            tables = tabula.read_pdf(
                pdf_path,
                pages='all',
                multiple_tables=True,
                pandas_options={'header': None},
                lattice=True  # helps detect tables better
            )
        except Exception as e:
            raise Exception(f"Tabula failed: {str(e)}")

        if not tables or len(tables) == 0:
            raise Exception("No tables found in this PDF. Try another file with proper tables.")

        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            for i, table in enumerate(tables):
                table.to_excel(writer, sheet_name=f"Table_{i+1}", index=False)

        return output_filename, output_path


    def convert_pdf_to_ppt(pdf_path):
        output_filename = f"converted_{uuid.uuid4().hex}.pptx"
        output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], output_filename)

        image_folder = os.path.join(current_app.config['OUTPUT_FOLDER'], f"ppt_pages_{uuid.uuid4().hex}")
        os.makedirs(image_folder, exist_ok=True)

        # IMPORTANT: Change this if your Poppler path is different
        poppler_path = r"C:\poppler-25.12.0\Library\bin"

        pages = convert_from_path(pdf_path, dpi=200, poppler_path=poppler_path)

        prs = Presentation()

        for i, page in enumerate(pages):
            img_path = os.path.join(image_folder, f'page_{i+1}.png')
            page.save(img_path, "PNG")

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            slide.shapes.add_picture(img_path, 0, 0, width=slide_width, height=slide_height)

        prs.save(output_path)

        return output_filename, output_path
    
    def convert_office_to_pdf(input_path):
        import subprocess
        
        ext = input_path.rsplit('.', 1)[1].lower()
        output_filename = f"converted_{uuid.uuid4().hex}.pdf"
        output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], output_filename)
        
        if ext not in ['doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx']:
            raise Exception("Unsupported office file type.")
        
        try:
            # Try using LibreOffice in headless mode (most reliable)
            abs_input = os.path.abspath(input_path)
            abs_output_dir = os.path.abspath(current_app.config['OUTPUT_FOLDER'])
            
            cmd = [
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', abs_output_dir,
                abs_input
            ]
            
            result = subprocess.run(cmd, capture_output=True, timeout=60, text=True)
            
            if result.returncode != 0:
                raise Exception(f"LibreOffice conversion failed: {result.stderr}")
            
            # LibreOffice outputs with the same base filename, just .pdf extension
            base_name = os.path.splitext(os.path.basename(input_path))[0]
            converted_path = os.path.join(abs_output_dir, f"{base_name}.pdf")
            
            if not os.path.exists(converted_path):
                raise Exception("PDF file was not created by LibreOffice")
            
            # Rename to our naming convention
            os.rename(converted_path, output_path)
            
        except Exception as e:
            # Fallback: Try using WIN32COM
            try:
                pythoncom.CoInitialize()
                
                if ext in ['doc', 'docx']:
                    word = win32com.client.Dispatch("Word.Application")
                    word.Visible = False
                    doc = word.Documents.Open(os.path.abspath(input_path))
                    doc.SaveAs(os.path.abspath(output_path), FileFormat=17)  # PDF
                    doc.Close()
                    word.Quit()

                elif ext in ['xls', 'xlsx']:
                    excel = win32com.client.Dispatch("Excel.Application")
                    excel.Visible = False
                    workbook = excel.Workbooks.Open(os.path.abspath(input_path))
                    workbook.ExportAsFixedFormat(0, os.path.abspath(output_path))  # 0 = PDF
                    workbook.Close(False)
                    excel.Quit()

                elif ext in ['ppt', 'pptx']:
                    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                    presentation = powerpoint.Presentations.Open(os.path.abspath(input_path), WithWindow=False)
                    presentation.SaveAs(os.path.abspath(output_path), 32)  # 32 = PDF
                    presentation.Close()
                    powerpoint.Quit()
                
                pythoncom.CoUninitialize()
                
            except Exception as fallback_error:
                raise Exception(f"All conversion methods failed - LibreOffice: {str(e)} | WIN32COM: {str(fallback_error)}")

        return output_filename, output_path

    # ===================== HOME =========================
    @app.route('/')
    def home():
        return render_template('index.html')
    
    # ===================== MERGE PDF =====================
    @app.route('/output/<filename>')
    def output_file(filename):
        return send_from_directory(app.config['OUTPUT_FOLDER'], filename, as_attachment=True)
    

    # ===================== MERGE PDF =====================
    @app.route('/merge', methods=['GET', 'POST'])
    def merge_pdf():
        if request.method == 'POST':
            files = request.files.getlist('pdfs')
            merger = PdfMerger()

            for file in files:
                filename, filepath = save_uploaded_file(file, 'pdf')
                if filename:
                    merger.append(filepath)

            output_filename = f"merged_{uuid.uuid4().hex}.pdf"
            output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], output_filename)
            merger.write(output_path)
            merger.close()

            file_size = os.path.getsize(output_path)
            file_size_kb = round(file_size / 1024, 2)

            return jsonify({
                "success": True,
                "filename": output_filename,
                "original_name": "merged.pdf",
                "file_size": file_size,
                "file_size_kb": file_size_kb,
                "message": "Your Pdfs are Merged"
            })

        return render_template('merge.html')
    


    # ===================== QR GENERATOR =====================
    @app.route('/qr_generator')
    def qr_generator():
        return render_template('qr_generator.html')
    

    # ===================== MERGE PD =====================
    @app.route('/generate_qr', methods=['post'])
    def generate_qr():
        data = request.get_json()
        url = data.get('url', '').strip()
        theme = data.get('theme', 'classic')

        if not url:
            return jsonify({'success': False, 'error': 'No URL Provided'})
        theme_map = {
            "classic":  {"fill": "black",   "back": "white"},
            "neon":     {"fill": "#8b5cf6", "back": "#0f172a"},
            "ocean":    {"fill": "#0891b2", "back": "#ecfeff"},
            "sunset":   {"fill": "#ea580c", "back": "#fff7ed"},
            "minimal":  {"fill": "#111827", "back": "#f9fafb"},
            "midnight": {"fill": "#22d3ee", "back": "#020617"},
        }

        colors = theme_map.get(theme, theme_map["classic"])

        filename = f"qr_{uuid.uuid4().hex}.png"
        output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], filename)

        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=12,
            border=4
        )
        qr.add_data(url)
        qr.make(fit=True)

        img = qr.make_image(fill_color=colors["fill"], back_color=colors["back"])
        img.save(output_path)

        return jsonify({
            'success': True, 
            'qr_url': f"/output/{filename}"
        })

    # ===================== SPLIT PDF =====================
    @app.route('/split', methods=['GET', 'POST'])
    def split_pdf():
        if request.method == 'POST':
            file = request.files.get('pdf')
            page_number_str = request.form.get('page_number', '')
            
            if not file or file.filename == '':
                return jsonify({'error': 'No file uploaded'}), 400
            
            if not page_number_str:
                return jsonify({'error': 'No page number specified'}), 400

            filename, filepath = save_uploaded_file(file, 'pdf')
            if not filename:
                return jsonify({'error': 'Invalid PDF file'}), 400

            try:
                page_number = int(page_number_str)
                reader = PdfReader(filepath)

                if page_number < 1 or page_number > len(reader.pages):
                    return jsonify({'error': f'Invalid page number. This PDF has only {len(reader.pages)} pages.'}), 400

                writer = PdfWriter()
                writer.add_page(reader.pages[page_number - 1])

                output_filename = f"split_page_{uuid.uuid4().hex}.pdf"
                output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], output_filename)
                with open(output_path, 'wb') as output_file:
                    writer.write(output_file)

                # Get file size
                file_size = os.path.getsize(output_path)

                return jsonify({
                    'success': True,
                    'filename': output_filename,
                    'page_number': page_number,
                    'filesize': file_size,
                    'filesize_formatted': format_file_size(file_size)
                })

            except ValueError:
                return jsonify({'error': 'Invalid page number'}), 400
            except Exception as e:
                return jsonify({'error': f'Error splitting PDF: {str(e)}'}), 500

        return render_template('split.html')
    
    @app.route('/download_split_pdf/<filename>', methods=['POST'])
    def download_split_pdf(filename):
        try:
            new_filename = request.json.get('new_filename', filename)
            new_filename = secure_filename(new_filename)
            if not new_filename.endswith('.pdf'):
                new_filename += '.pdf'
            
            filepath = os.path.join(current_app.config['OUTPUT_FOLDER'], filename)
            
            if not os.path.exists(filepath) or not os.path.isfile(filepath):
                return jsonify({'error': 'File not found'}), 404
            
            return send_file(filepath, as_attachment=True, download_name=new_filename)
        except Exception as e:
            return jsonify({'error': str(e)}), 500

    # ===================== ROTATE PDF =====================
    @app.route('/rotate', methods=['GET', 'POST'])
    def rotate_pdf():
        if request.method == 'POST':
            file = request.files['pdf']
            angle = int(request.form['angle'])

            filename, filepath = save_uploaded_file(file, 'pdf')
            if not filename:
                return "<h2>No file uploaded or invalid PDF.</h2>"

            reader = PdfReader(filepath)
            writer = PdfWriter()

            for page in reader.pages:
                page.rotate_clockwise(angle)
                writer.add_page(page)

            output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], 'rotated.pdf')
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)

            return send_file(output_path, as_attachment=True)

        return render_template('rotate.html')

    # ===================== ORGANIZE PDF =====================
    @app.route('/organize', methods=['GET', 'POST'])
    def organize_pdf():
        if request.method == 'POST':
            file = request.files['pdf']
            pages_to_keep = request.form['pages']

            filename, filepath = save_uploaded_file(file, 'pdf')
            if not filename:
                return "<h2>No file uploaded or invalid PDF.</h2>"

            reader = PdfReader(filepath)
            writer = PdfWriter()

            try:
                pages = [int(p.strip()) - 1 for p in pages_to_keep.split(',')]
            except ValueError:
                return "<h2>Invalid input. Please enter pages like: 1,3,5</h2>"

            for p in pages:
                if 0 <= p < len(reader.pages):
                    writer.add_page(reader.pages[p])

            output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], 'organized.pdf')
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)

            return send_file(output_path, as_attachment=True)

        return render_template('organize.html')
    
    # ===================== COMPRESS IMAGE =====================
    @app.route('/compress_image', methods=['GET', 'POST'])
    def compress_image_tool():
        if request.method == 'POST':
            file = request.files.get('image')

            if not file:
                return jsonify({"success": False, "error": "No image uploaded"}), 400

            # Save input file
            filename = secure_filename(file.filename)
            input_path = os.path.join(current_app.config['UPLOAD_FOLDER'], filename)
            file.save(input_path)

            # Output path
            output_filename = f"compressed_{uuid.uuid4().hex}.jpg"
            output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], output_filename)

            # Get form values
            mode = request.form.get('mode')
            target_size = request.form.get('target_size')

            # Open image
            img = Image.open(input_path)

            # Convert to RGB (important for JPG)
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")

            # 🎯 Compression logic

            # If user gives target size
            if target_size:
                target_size = int(target_size) * 1024  # KB → bytes
                quality = 90

                while quality > 10:
                    img.save(output_path, "JPEG", quality=quality, optimize=True)
                    if os.path.getsize(output_path) <= target_size:
                        break
                    quality -= 5

            else:
                # Mode-based compression
                if mode == "maximum":
                    quality = 30
                elif mode == "medium":
                    quality = 60
                else:  # minimum
                    quality = 85

                img.save(output_path, "JPEG", quality=quality, optimize=True)

            file_size = os.path.getsize(output_path)
            file_size_kb = round(file_size / 1024, 2)

            return jsonify({
                "success": True,
                "filename": output_filename,
                "original_name": "compressed_image.jpg",
                "file_size": file_size,
                "file_size_kb": file_size_kb,
                "message": "Your fille is Image is Compressed"
            })

        return render_template('compress_image.html')

    # ===================== JPG TO PDF =====================
    @app.route('/jpg_to_pdf', methods=['GET', 'POST'])
    def jpg_to_pdf():
        if request.method == 'POST':
            files = request.files.getlist('images')
            cleaned_image_paths = []

            for index, file in enumerate(files):
                filename, filepath = save_uploaded_file(file, 'image')
                if not filename:
                    continue

                try:
                    img = Image.open(filepath)
                    img = ImageOps.exif_transpose(img)

                    if img.mode in ("RGBA", "P"):
                        img = img.convert("RGB")

                    cleaned_path = os.path.join(current_app.config['UPLOAD_FOLDER'], f"cleaned_{index}.jpg")
                    img.save(cleaned_path, "JPEG")
                    cleaned_image_paths.append(cleaned_path)

                except Exception as e:
                    return jsonify({"success": False, "error": f"Error processing image {filename}: {str(e)}"})

            if not cleaned_image_paths:
                return jsonify({"success": False, "error": "No valid image files uploaded."})

            output_filename = f"converted_{uuid.uuid4().hex}.pdf"
            output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], output_filename)

            try:
                with open(output_path, "wb") as f:
                    f.write(img2pdf.convert(cleaned_image_paths))
            except Exception as e:
                return jsonify({"success": False, "error": f"Error converting images to PDF: {str(e)}"})

            # Get file size
            file_size = os.path.getsize(output_path)
            file_size_kb = round(file_size / 1024, 2)
            
            # Clean up temporary files
            for cleaned_path in cleaned_image_paths:
                try:
                    os.remove(cleaned_path)
                except:
                    pass

            return jsonify({
                "success": True,
                "filename": output_filename,
                "original_name": "converted.pdf",
                "file_size": file_size,
                "file_size_kb": file_size_kb,
                "message": "Your file is converted!"
            })

        return render_template('jpg_to_pdf.html')

    # ===================== PDF TO JPG (NO POPPLER) =====================
    @app.route('/pdf_to_jpg', methods=['GET', 'POST'])
    def pdf_to_jpg():
        if request.method == 'POST':
            file = request.files.get('pdf')
            filename, filepath = save_uploaded_file(file, 'pdf')
            if not filename:
                return jsonify({"success": False, "error": "No file uploaded or invalid PDF."})

            import tempfile
            import shutil
            try:
                doc = fitz.open(filepath)
                # Generate unique ID for output
                output_filename = f"converted_images_{uuid.uuid4().hex}.zip"
                final_zip_path = os.path.join(current_app.config['OUTPUT_FOLDER'], output_filename)

                with tempfile.TemporaryDirectory() as temp_dir:
                    temp_zip_path = os.path.join(temp_dir, output_filename)
                    
                    with zipfile.ZipFile(temp_zip_path, 'w') as zipf:
                        for i, page in enumerate(doc):
                            pix = page.get_pixmap()
                            img_path = os.path.join(temp_dir, f'page_{i+1}.jpg')
                            pix.save(img_path)
                            zipf.write(img_path, os.path.basename(img_path))
                            
                    # Atomically move the completed zip to the server output directory
                    shutil.move(temp_zip_path, final_zip_path)

                doc.close()
                
                file_size = os.path.getsize(final_zip_path)
                file_size_kb = round(file_size / 1024, 2)

                return jsonify({
                    "success": True,
                    "filename": output_filename,
                    "original_name": "converted_images.zip",
                    "file_size": file_size,
                    "file_size_kb": file_size_kb,
                    "message": "Your Pdfs is Converted into the Image"
                })

            except Exception as e:
                return jsonify({"success": False, "error": f"Error converting PDF to JPG: {str(e)}"})

        return render_template('pdf_to_jpg.html')

    # ===================== PDF TO TEXT =====================
    @app.route('/pdf_to_text', methods=['GET', 'POST'])
    def pdf_to_text():
        if request.method == 'POST':
            file = request.files.get('pdf')
            pages_input = request.form.get('pages', '').strip()
            
            filename, filepath = save_uploaded_file(file, 'pdf')
            if not filename:
                return jsonify({"success": False, "error": "No file uploaded or invalid PDF."})

            try:
                text_output = ""

                with pdfplumber.open(filepath) as pdf:
                    total_pages = len(pdf.pages)
                    pages_to_extract = []

                    if pages_input:
                        try:
                            # Parse page numbers like "1, 3, 5-7"
                            for part in pages_input.split(','):
                                part = part.strip()
                                if '-' in part:
                                    start, end = map(int, part.split('-'))
                                    pages_to_extract.extend(range(start, end + 1))
                                else:
                                    pages_to_extract.append(int(part))
                            
                            # Filter and adjust to 0-indexed
                            pages_to_extract = [p - 1 for p in pages_to_extract if 1 <= p <= total_pages]
                        except ValueError:
                            return jsonify({"success": False, "error": "Invalid page numbers format."})
                    else:
                        pages_to_extract = list(range(total_pages))

                    if not pages_to_extract:
                        return jsonify({"success": False, "error": "No valid pages selected for extraction. Check page limit."})

                    for page_num in pages_to_extract:
                        text = pdf.pages[page_num].extract_text()
                        if text:
                            if len(pages_to_extract) > 1:
                                text_output += f"--- Page {page_num + 1} ---\n{text}\n\n"
                            else:
                                text_output += f"{text}\n\n"

                output_filename = f"extracted_text_{uuid.uuid4().hex}.txt"
                output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], output_filename)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text_output)

                file_size = os.path.getsize(output_path)
                file_size_kb = round(file_size / 1024, 2)

                return jsonify({
                    "success": True,
                    "filename": output_filename,
                    "original_name": "extracted_text.txt",
                    "file_size": file_size,
                    "file_size_kb": file_size_kb,
                    "message": "Your Text file is Converted"
                })

            except Exception as e:
                return jsonify({"success": False, "error": f"Error converting PDF to Text: {str(e)}"})

        return render_template('pdf_to_text.html')
    
    @app.route('/download_text_file/<filename>', methods=['POST'])
    def download_text_file(filename):
        try:
            # Get new filename if user renamed it
            new_filename = request.json.get('new_filename', filename)
            new_filename = secure_filename(new_filename)
            if not new_filename.endswith('.txt'):
                new_filename += '.txt'
            
            filepath = os.path.join(current_app.config['OUTPUT_FOLDER'], filename)
            
            # Security check: ensure file exists in output folder
            if not os.path.exists(filepath) or not os.path.isfile(filepath):
                return jsonify({'error': 'File not found'}), 404
            
            return send_file(filepath, as_attachment=True, download_name=new_filename)
        except Exception as e:
            return jsonify({'error': str(e)}), 500
    
    def format_file_size(bytes):
        for unit in ['B', 'KB', 'MB', 'GB']:
            if bytes < 1024:
                return f"{bytes:.2f} {unit}"
            bytes /= 1024
        return f"{bytes:.2f} TB"

    # ===================== PDF EXTRACT PAGES =====================
    @app.route('/pdf_extract_pages', methods=['GET', 'POST'])
    def pdf_extract_pages():
        if request.method == 'POST':
            file = request.files.get('pdf')
            page_numbers = request.form.get('pages', '').strip()
            
            if not file or file.filename == '':
                return jsonify({'error': 'No file uploaded'}), 400
            
            if not page_numbers:
                return jsonify({'error': 'No page numbers specified'}), 400
            
            filename, filepath = save_uploaded_file(file, 'pdf')
            if not filename:
                return jsonify({'error': 'Invalid PDF file'}), 400

            try:
                # Parse page numbers (supports "1,3,5" or "1-5" or "1")
                pages_to_extract = parse_page_numbers(page_numbers, filepath)
                
                if not pages_to_extract:
                    return jsonify({'error': 'No valid page numbers found'}), 400
                
                # Extract pages
                output_filename = f"extracted_pages_{uuid.uuid4().hex}.pdf"
                output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], output_filename)
                
                pdf_reader = PdfReader(filepath)
                pdf_writer = PdfWriter()
                
                for page_num in pages_to_extract:
                    if 0 <= page_num < len(pdf_reader.pages):
                        pdf_writer.add_page(pdf_reader.pages[page_num])
                
                with open(output_path, 'wb') as output_file:
                    pdf_writer.write(output_file)
                
                # Get file size
                file_size = os.path.getsize(output_path)
                
                return jsonify({
                    'success': True,
                    'filename': output_filename,
                    'pages': page_numbers,
                    'filesize': file_size,
                    'filesize_formatted': format_file_size(file_size)
                })
            
            except Exception as e:
                return jsonify({'error': f'Error extracting pages: {str(e)}'}), 500
        
        return render_template('pdf_extract_pages.html')
    
    def parse_page_numbers(page_string, filepath):
        """Parse page numbers from string like '1,3,5' or '1-5' or '1'"""
        pages = set()
        
        # Get total pages in PDF
        try:
            with pdfplumber.open(filepath) as pdf:
                total_pages = len(pdf.pages)
        except:
            return []
        
        parts = page_string.split(',')
        for part in parts:
            part = part.strip()
            if '-' in part:
                # Range like "1-5"
                try:
                    start, end = part.split('-')
                    start = int(start.strip()) - 1  # Convert to 0-based indexing
                    end = int(end.strip())  # Keep as 1-based for range
                    for i in range(start, min(end, total_pages)):
                        pages.add(i)
                except:
                    pass
            else:
                # Single page
                try:
                    page_num = int(part) - 1  # Convert to 0-based indexing
                    if 0 <= page_num < total_pages:
                        pages.add(page_num)
                except:
                    pass
        
        return sorted(list(pages))
    
    @app.route('/download_extracted_pdf/<filename>', methods=['POST'])
    def download_extracted_pdf(filename):
        try:
            new_filename = request.json.get('new_filename', filename)
            new_filename = secure_filename(new_filename)
            if not new_filename.endswith('.pdf'):
                new_filename += '.pdf'
            
            filepath = os.path.join(current_app.config['OUTPUT_FOLDER'], filename)
            
            if not os.path.exists(filepath) or not os.path.isfile(filepath):
                return jsonify({'error': 'File not found'}), 404
            
            return send_file(filepath, as_attachment=True, download_name=new_filename)
        except Exception as e:
            return jsonify({'error': str(e)}), 500

    # ===================== HTML TO PDF =====================
    @app.route('/html_to_pdf', methods=['GET', 'POST'])
    def html_to_pdf():
        if request.method == 'POST':
            html_content = request.form['html_content']

            output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], 'html_to_pdf.pdf')

            c = canvas.Canvas(output_path, pagesize=letter)
            width, height = letter

            y = height - 40
            lines = html_content.split('\n')

            c.setFont("Helvetica", 10)

            for line in lines:
                if y < 40:
                    c.showPage()
                    c.setFont("Helvetica", 10)
                    y = height - 40
                c.drawString(40, y, line[:100])
                y -= 15

            c.save()

            return send_file(output_path, as_attachment=True)

        return render_template('html_to_pdf.html')
    
        # ===================== PDF CONVERTER =====================
    @app.route('/pdf_converter', methods=['GET', 'POST'])
    def pdf_converter():
        if request.method == 'POST':
            file = request.files.get('file')
            convert_to = request.form.get('convert_to')

            filename, filepath = save_uploaded_file(file, 'pdf')
            if not filename:
                return "<h2>No file uploaded or invalid PDF.</h2>"

            if not convert_to:
                return "<h2>Please select a conversion format.</h2>"

            try:
                if convert_to == 'word':
                    output_filename, output_path = convert_pdf_to_word(filepath)

                elif convert_to == 'excel':
                    output_filename, output_path = convert_pdf_to_excel(filepath)

                elif convert_to == 'ppt':
                    output_filename, output_path = convert_pdf_to_ppt(filepath)

                else:
                    return "<h2>Invalid conversion format selected.</h2>"

                return send_file(output_path, as_attachment=True)

            except Exception as e:
                return f"<h2>Conversion failed: {str(e)}</h2>"

        return render_template('pdf_converter.html')
    
    # ===================== OFFICE TO PDF =====================

    @app.route('/office_to_pdf', methods=['GET', 'POST'])
    def office_to_pdf():
        if request.method == 'POST':
            file = request.files.get('file')

            if not file or file.filename == '':
                return jsonify({'error': 'No file uploaded'}), 400

            if not allowed_office_file(file.filename):
                return jsonify({'error': 'Invalid file type. Please upload Word, Excel, or PowerPoint file'}), 400

            filename = secure_filename(file.filename)
            filepath = os.path.join(current_app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            try:
                output_filename, output_path = convert_office_to_pdf(filepath)
                
                # Return the PDF file directly
                return send_file(
                    output_path,
                    mimetype='application/pdf',
                    as_attachment=True,
                    download_name='converted.pdf'
                )

            except Exception as e:
                return jsonify({'error': f'Conversion failed: {str(e)}'}), 500

        return render_template('office_to_pdf.html')
    
    @app.route('/download_office_pdf/<filename>', methods=['POST'])
    def download_office_pdf(filename):
        try:
            new_filename = request.json.get('new_filename', filename)
            new_filename = secure_filename(new_filename)
            if not new_filename.endswith('.pdf'):
                new_filename += '.pdf'
            
            filepath = os.path.join(current_app.config['OUTPUT_FOLDER'], filename)
            
            if not os.path.exists(filepath) or not os.path.isfile(filepath):
                return jsonify({'error': 'File not found'}), 404
            
            return send_file(filepath, as_attachment=True, download_name=new_filename)
        except Exception as e:
            return jsonify({'error': str(e)}), 500
    

# ===================== PAGES =====================
    @app.route('/Privacypolicy')
    def privacy_policy():
        return render_template('Privacypolicy.html')


    @app.route('/Terms&Condition')
    def terms_condition():
        return render_template('Terms&Condition.html')


    @app.route('/about')
    def about():
        return render_template('about.html')


    @app.route('/Contactus')
    def contact():
        return render_template('Contactus.html')
    
    # ===================== DOWNLOAD =====================
    @app.route('/download/<filename>')
    def download_file(filename):
        """Download the converted file with an optional custom name"""
        # Get custom filename from query parameter, default to original filename
        custom_name = request.args.get('name', filename)
        custom_name = secure_filename(custom_name)
        
        if not custom_name or not filename:
            return jsonify({"error": "Invalid filename"}), 400
        
        file_path = os.path.join(current_app.config['OUTPUT_FOLDER'], filename)
        
        # Security check: ensure file exists and is in outputs folder
        if not os.path.exists(file_path) or not os.path.isfile(file_path):
            return jsonify({"error": "File not found"}), 404
        
        return send_file(file_path, as_attachment=True, download_name=custom_name)

    # ===================== COMPRESS PDF =====================
    @app.route('/compress', methods=['GET', 'POST'])
    def compress_pdf():
        if request.method == 'POST':
            file = request.files['pdf']

            filename, filepath = save_uploaded_file(file, 'pdf')
            if not filename:
                return "<h2>No file uploaded or invalid PDF.</h2>"

            output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], 'compressed.pdf')

            # Original size
            original_size = os.path.getsize(filepath)

            # Compress using PyMuPDF
            doc = fitz.open(filepath)
            doc.save(output_path, garbage=4, deflate=True)
            doc.close()

            # Compressed size
            compressed_size = os.path.getsize(output_path)

            # Reduction calculations
            reduced_bytes = original_size - compressed_size
            reduction_percent = (reduced_bytes / original_size) * 100 if original_size > 0 else 0

            return render_template(
                'compress_result.html',
                original_size=round(original_size / 1024, 2),
                compressed_size=round(compressed_size / 1024, 2),
                reduced_size=round(reduced_bytes / 1024, 2),
                reduction_percent=round(reduction_percent, 2),
                download_file='compressed.pdf'
          )

        return render_template('compress.html')
    
    return app

if __name__ == '__main__':
    app = create_app()
    app.run(debug=True)